"""
ppt_generator.py — Premium McKinsey-style PowerPoint intelligence report.
Dark consulting theme, gradient accents, metric cards, heatmap slide.
"""

import io
import math
import statistics
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Color Palette ────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x07, 0x0B, 0x14)
BG_CARD       = RGBColor(0x0D, 0x15, 0x26)
BG_CARD2      = RGBColor(0x11, 0x1A, 0x2E)
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)
ACCENT_VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_PINK   = RGBColor(0xEC, 0x48, 0x99)
ACCENT_CYAN   = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_GREEN  = RGBColor(0x10, 0xB9, 0x81)
ACCENT_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED    = RGBColor(0xEF, 0x44, 0x44)
TEXT_PRIMARY  = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_MUTED    = RGBColor(0x94, 0xA3, 0xB8)
TEXT_DIM      = RGBColor(0x47, 0x55, 0x69)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER       = RGBColor(0x1E, 0x29, 0x3B)

COLORS_RGB = [ACCENT_PURPLE, ACCENT_PINK, ACCENT_CYAN, ACCENT_GREEN, ACCENT_AMBER]


def fmt_num(n):
    if n >= 1_000_000:
        return f"{n/1_000_000:.1f}M"
    elif n >= 1_000:
        return f"{n/1_000:.1f}K"
    return str(int(n))


def _heatmap_color(score, low=0, high=10):
    t = max(0.0, min(1.0, (score - low) / max(high - low, 1)))
    if t < 0.5:
        r = int(0x7F + (0x6D - 0x7F) * t * 2)
        g = int(0x1D + (0x28 - 0x1D) * t * 2)
        b = int(0x1D + (0xD9 - 0x1D) * t * 2)
    else:
        t2 = (t - 0.5) * 2
        r = int(0x6D + (0x10 - 0x6D) * t2)
        g = int(0x28 + (0xB9 - 0x28) * t2)
        b = int(0xD9 + (0x81 - 0xD9) * t2)
    return RGBColor(r, g, b)


# ─── Core Helpers ─────────────────────────────────────────────────────────────

def _add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_DARK
    return slide


def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _text(slide, txt, l, t, w, h, size=11, bold=False, color=TEXT_PRIMARY,
          align=PP_ALIGN.LEFT, italic=False):
    if not txt:
        return
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(txt)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def _header(slide, prs, title, subtitle="", slide_num=None):
    W = prs.slide_width
    _rect(slide, 0, 0, W, Inches(0.06), ACCENT_PURPLE)
    _rect(slide, 0, Inches(0.06), W, Inches(0.95), BG_CARD)
    _text(slide, title, Inches(0.45), Inches(0.1), W - Inches(2), Inches(0.55),
          size=22, bold=True, color=TEXT_PRIMARY)
    if subtitle:
        _text(slide, subtitle, Inches(0.45), Inches(0.62), W - Inches(2), Inches(0.32),
              size=10, color=TEXT_MUTED)
    if slide_num:
        _rect(slide, W - Inches(0.8), Inches(0.2), Inches(0.55), Inches(0.32), ACCENT_PURPLE)
        _text(slide, str(slide_num), W - Inches(0.8), Inches(0.2), Inches(0.55), Inches(0.32),
              size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, 0, Inches(1.01), W, Inches(0.02), DIVIDER)


def _kpi_card(slide, label, value, sub, l, t, w=Inches(2.2), h=Inches(1.15), accent=ACCENT_PURPLE):
    _rect(slide, l, t, w, h, BG_CARD2)
    _rect(slide, l, t, Inches(0.05), h, accent)
    _rect(slide, l, t + h - Inches(0.03), w, Inches(0.03), accent)
    _text(slide, label.upper(), l + Inches(0.12), t + Inches(0.1), w - Inches(0.2), Inches(0.22),
          size=7, bold=True, color=TEXT_MUTED)
    _text(slide, value, l + Inches(0.12), t + Inches(0.3), w - Inches(0.2), Inches(0.52),
          size=20, bold=True, color=TEXT_PRIMARY)
    _text(slide, sub, l + Inches(0.12), t + Inches(0.85), w - Inches(0.2), Inches(0.25),
          size=8, color=accent, italic=True)


def _insight(slide, text, l, t, w, h, accent=ACCENT_PURPLE, icon="->"):
    _rect(slide, l, t, w, h, BG_CARD2)
    _rect(slide, l, t, Inches(0.05), h, accent)
    _text(slide, f"{icon}  {text}", l + Inches(0.12), t + Inches(0.1),
          w - Inches(0.18), h - Inches(0.15), size=9.5, color=TEXT_PRIMARY)


def _hbar(slide, companies, values, title, l, t, w, h, color=ACCENT_PURPLE, unit=""):
    _text(slide, title, l, t, w, Inches(0.28), size=10, bold=True, color=TEXT_MUTED)
    t += Inches(0.32)
    h -= Inches(0.32)
    if not values or max(values) == 0:
        return
    bar_h = h / max(len(companies), 1)
    max_val = max(values)
    bar_area_w = w - Inches(1.5)
    for i, (co, val) in enumerate(zip(companies, values)):
        by = t + i * bar_h + bar_h * 0.12
        bh = bar_h * 0.65
        _text(slide, co[:14], l, by, Inches(1.35), bh, size=8.5, color=TEXT_MUTED)
        _rect(slide, l + Inches(1.45), by + bh * 0.3, bar_area_w, bh * 0.35,
              RGBColor(0x1E, 0x29, 0x3B))
        bw = max((val / max_val) * bar_area_w, Inches(0.05))
        _rect(slide, l + Inches(1.45), by + bh * 0.3, bw, bh * 0.35,
              COLORS_RGB[i % len(COLORS_RGB)])
        _text(slide, f"{fmt_num(val)}{unit}",
              l + Inches(1.45) + bw + Inches(0.06), by,
              Inches(0.9), bh, size=8.5, bold=True, color=TEXT_PRIMARY)


def _section_label(slide, text, l, t, w=Inches(3)):
    _rect(slide, l, t, Inches(0.03), Inches(0.22), ACCENT_PURPLE)
    _text(slide, text.upper(), l + Inches(0.08), t, w, Inches(0.22),
          size=7.5, bold=True, color=ACCENT_VIOLET)


# ─── Slide Builders ───────────────────────────────────────────────────────────

def _slide_title(prs, user_company, companies):
    slide = _add_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    _rect(slide, 0, 0, W, H, RGBColor(0x05, 0x08, 0x12))
    _rect(slide, 0, 0, Inches(0.12), H, ACCENT_PURPLE)
    _rect(slide, Inches(0.12), 0, Inches(0.04), H, RGBColor(0x4C, 0x1D, 0x95))
    _rect(slide, W - Inches(3.8), 0, Inches(3.8), H, RGBColor(0x0D, 0x10, 0x1F))
    _rect(slide, W - Inches(3.8), 0, Inches(0.02), H, RGBColor(0x1E, 0x29, 0x3B))
    _rect(slide, Inches(0.6), Inches(1.2), Inches(2.8), Inches(0.28), RGBColor(0x1E, 0x10, 0x3C))
    _text(slide, "COMPETITIVE INTELLIGENCE PLATFORM", Inches(0.65), Inches(1.22),
          Inches(2.8), Inches(0.26), size=7, bold=True, color=ACCENT_VIOLET)
    _text(slide, "AI VIDEO", Inches(0.6), Inches(1.65), Inches(8), Inches(0.9),
          size=46, bold=True, color=TEXT_PRIMARY)
    _text(slide, "MARKETING", Inches(0.6), Inches(2.45), Inches(8), Inches(0.9),
          size=46, bold=True, color=ACCENT_VIOLET)
    _text(slide, "WAR ROOM", Inches(0.6), Inches(3.25), Inches(8), Inches(0.9),
          size=46, bold=True, color=ACCENT_PINK)
    _rect(slide, Inches(0.6), Inches(4.25), Inches(5.5), Inches(0.03), ACCENT_PURPLE)
    _text(slide, f"Subject Company:  {user_company}", Inches(0.6), Inches(4.4),
          Inches(6), Inches(0.3), size=11, color=TEXT_MUTED)
    rivals = [c for c in companies if c != user_company]
    _text(slide, f"Competitors:  {' . '.join(rivals)}", Inches(0.6), Inches(4.75),
          Inches(6), Inches(0.3), size=10, color=TEXT_DIM)
    _text(slide, f"Generated  {datetime.now().strftime('%B %d, %Y')}",
          Inches(0.6), Inches(5.1), Inches(6), Inches(0.3), size=9, color=TEXT_DIM)
    _text(slide, "REPORT CONTENTS", W - Inches(3.5), Inches(1.0), Inches(3.2), Inches(0.28),
          size=7.5, bold=True, color=TEXT_DIM)
    contents = [
        "01  Executive Summary",
        "02  Channel Intelligence",
        "03  Engagement Analysis",
        "04  Attention Efficiency Index",
        "05  Content Theme Map",
        "06  Opportunity Heatmap",
        "07  Competitive Threat Matrix",
        "08  AI Strategic Recommendations",
        "09  30-Day Execution Roadmap",
        "10  Next Steps",
    ]
    for i, item in enumerate(contents):
        _text(slide, item, W - Inches(3.5), Inches(1.4) + i * Inches(0.55),
              Inches(3.2), Inches(0.45), size=9,
              color=ACCENT_VIOLET if i == 0 else TEXT_MUTED)
    _text(slide, "CONFIDENTIAL  .  STRATEGIC INTELLIGENCE",
          Inches(0.6), H - Inches(0.45), Inches(8), Inches(0.3),
          size=7.5, bold=True, color=TEXT_DIM)


def _slide_exec_summary(prs, analytics, user_company, slide_num):
    slide = _add_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    es = analytics.executive_summary
    metrics = analytics.metrics
    user_m = metrics.get(user_company, {})
    _header(slide, prs, "Executive Summary",
            "Strategic competitive positioning at a glance", slide_num)
    kpi_data = [
        ("Threat Score",    f"{es.get('user_threat_score',0)}/10",   "Overall power rating",    ACCENT_PURPLE),
        ("Top Rival",       es.get('strongest_rival','N/A'),          f"Score {es.get('strongest_rival_score',0)}/10", ACCENT_PINK),
        ("Subscribers",     fmt_num(user_m.get('subscribers',0)),     user_company,              ACCENT_CYAN),
        ("Engagement",      f"{user_m.get('engagement_rate',0):.1f}%","Avg across videos",       ACCENT_GREEN),
        ("Top Opportunity", es.get('top_opportunity','N/A'),          "White space theme",       ACCENT_AMBER),
    ]
    for i, (label, val, sub, acc) in enumerate(kpi_data):
        _kpi_card(slide, label, val, sub,
                  Inches(0.3) + i * Inches(2.6), Inches(1.15),
                  Inches(2.4), Inches(1.2), acc)
    _section_label(slide, "Strategic Intelligence Findings", Inches(0.3), Inches(2.55))
    insights = [
        (es.get("insight_1",""), ACCENT_PINK,  "!"),
        (es.get("insight_2",""), ACCENT_GREEN, "V"),
        (es.get("insight_3",""), ACCENT_AMBER, ">"),
    ]
    for i, (text, acc, icon) in enumerate(insights):
        _insight(slide, text, Inches(0.3), Inches(2.82) + i * Inches(0.95),
                 W - Inches(0.6), Inches(0.85), acc, icon)
    _rect(slide, 0, H - Inches(0.35), W, Inches(0.35), BG_CARD)
    _text(slide, f"Analysis covers {len(metrics)} companies  .  Data: YouTube API v3  .  {datetime.now().strftime('%B %Y')}",
          Inches(0.4), H - Inches(0.3), W - Inches(0.8), Inches(0.25),
          size=7.5, color=TEXT_DIM)


def _slide_channel_intel(prs, analytics, channel_data, slide_num):
    slide = _add_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    metrics = analytics.metrics
    companies = list(metrics.keys())
    _header(slide, prs, "Channel Intelligence",
            "Subscriber base, publishing volume, and upload velocity", slide_num)
    card_w = W / max(len(companies), 1) - Inches(0.15)
    for i, c in enumerate(companies):
        m = metrics[c]
        acc = COLORS_RGB[i % len(COLORS_RGB)]
        cx = Inches(0.1) + i * (card_w + Inches(0.1))
        _rect(slide, cx, Inches(1.1), card_w, Inches(0.9), BG_CARD2)
        _rect(slide, cx, Inches(1.1), card_w, Inches(0.04), acc)
        _text(slide, c[:14], cx + Inches(0.1), Inches(1.15), card_w - Inches(0.2), Inches(0.25),
              size=8.5, bold=True, color=acc)
        _text(slide, fmt_num(m["subscribers"]), cx + Inches(0.1), Inches(1.38),
              card_w - Inches(0.2), Inches(0.4), size=18, bold=True, color=TEXT_PRIMARY)
        _text(slide, "subscribers", cx + Inches(0.1), Inches(1.75),
              card_w - Inches(0.2), Inches(0.2), size=7.5, color=TEXT_MUTED)
    col_w = (W - Inches(0.8)) / 2
    subs = [metrics[c]["subscribers"] for c in companies]
    _hbar(slide, companies, subs, "SUBSCRIBER COUNT",
          Inches(0.3), Inches(2.15), col_w, Inches(2.1))
    freqs = [metrics[c]["upload_freq"] for c in companies]
    _hbar(slide, companies, freqs, "UPLOAD FREQUENCY (videos/month)",
          Inches(0.4) + col_w, Inches(2.15), col_w, Inches(2.1), color=ACCENT_CYAN, unit="/mo")
    total_vids = [metrics[c]["total_videos"] for c in companies]
    _hbar(slide, companies, total_vids, "TOTAL VIDEOS PUBLISHED",
          Inches(0.3), Inches(4.4), col_w, Inches(2.1), color=ACCENT_PINK)
    total_views = [metrics[c]["total_views"] for c in companies]
    _hbar(slide, companies, total_views, "TOTAL CHANNEL VIEWS",
          Inches(0.4) + col_w, Inches(4.4), col_w, Inches(2.1), color=ACCENT_GREEN)


def _slide_engagement(prs, analytics, slide_num):
    slide = _add_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    metrics = analytics.metrics
    companies = list(metrics.keys())
    _header(slide, prs, "Engagement Intelligence",
            "Views, likes, comments, and engagement rate analysis", slide_num)
    for i, c in enumerate(companies[:5]):
        m = metrics[c]
        acc = COLORS_RGB[i % len(COLORS_RGB)]
        _kpi_card(slide, c[:12], fmt_num(m["avg_views"]),
                  f"{m['engagement_rate']:.1f}% engagement",
                  Inches(0.3) + i * Inches(2.6), Inches(1.1),
                  Inches(2.4), Inches(1.1), acc)
    col_w = (W - Inches(0.8)) / 2
    eng = [metrics[c]["engagement_rate"] for c in companies]
    _hbar(slide, companies, eng, "ENGAGEMENT RATE (%)",
          Inches(0.3), Inches(2.4), col_w, Inches(2.0), color=ACCENT_PURPLE, unit="%")
    avg_views = [metrics[c]["avg_views"] for c in companies]
    _hbar(slide, companies, avg_views, "AVERAGE VIEWS PER VIDEO",
          Inches(0.4) + col_w, Inches(2.4), col_w, Inches(2.0), color=ACCENT_CYAN)
    avg_likes = [metrics[c]["avg_likes"] for c in companies]
    _hbar(slide, companies, avg_likes, "AVERAGE LIKES PER VIDEO",
          Inches(0.3), Inches(4.55), col_w, Inches(2.0), color=ACCENT_PINK)
    avg_comments = [metrics[c]["avg_comments"] for c in companies]
    _hbar(slide, companies, avg_comments, "AVERAGE COMMENTS PER VIDEO",
          Inches(0.4) + col_w, Inches(4.55), col_w, Inches(2.0), color=ACCENT_GREEN)


def _slide_aei(prs, analytics, slide_num):
    slide = _add_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    aei = analytics.aei_scores
    companies_sorted = sorted(aei.keys(), key=lambda c: aei[c], reverse=True)
    _header(slide, prs, "Attention Efficiency Index (AEI)",
            "How effectively each brand converts subscribers into active viewers", slide_num)
    _rect(slide, Inches(0.3), Inches(1.1), W - Inches(0.6), Inches(0.5), BG_CARD2)
    _rect(slide, Inches(0.3), Inches(1.1), Inches(0.05), Inches(0.5), ACCENT_PURPLE)
    _text(slide, "AEI = log-normalized(Avg Views / Subscribers x 100)   |   >60 = Excellent   30-60 = Healthy   <30 = Needs Activation",
          Inches(0.45), Inches(1.17), W - Inches(0.8), Inches(0.38),
          size=9, color=TEXT_MUTED, italic=True)
    tier_map = [(60, 100, "EXCELLENT", ACCENT_GREEN),
                (30, 60,  "HEALTHY",   ACCENT_AMBER),
                (0,  30,  "ACTIVATE",  ACCENT_RED)]
    for rank, c in enumerate(companies_sorted):
        score = aei[c]
        tier, acc = next(((t, a) for lo, hi, t, a in tier_map if lo <= score <= hi),
                         ("ACTIVATE", ACCENT_RED))
        medal = ["#1", "#2", "#3", "#4", "#5"][rank] if rank < 5 else str(rank+1)
        card_y = Inches(1.75) + rank * Inches(0.98)
        card_h = Inches(0.85)
        _rect(slide, Inches(0.3), card_y, W - Inches(0.6), card_h, BG_CARD2)
        _rect(slide, Inches(0.3), card_y, Inches(0.05), card_h, acc)
        _text(slide, f"{medal}  {c}", Inches(0.45), card_y + Inches(0.08),
              Inches(3), Inches(0.35), size=12, bold=True, color=TEXT_PRIMARY)
        _rect(slide, Inches(3.7), card_y + Inches(0.18), Inches(1.0), Inches(0.3), acc)
        _text(slide, tier, Inches(3.7), card_y + Inches(0.18), Inches(1.0), Inches(0.3),
              size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(slide, f"{score:.1f}", Inches(5.1), card_y + Inches(0.05),
              Inches(1.5), Inches(0.55), size=26, bold=True, color=acc)
        _text(slide, "/ 100", Inches(6.5), card_y + Inches(0.28),
              Inches(0.7), Inches(0.3), size=10, color=TEXT_MUTED)
        bar_start = Inches(7.5)
        bar_w = W - Inches(8.0)
        _rect(slide, bar_start, card_y + Inches(0.3), bar_w, Inches(0.2),
              RGBColor(0x1E, 0x29, 0x3B))
        fill_w = max(bar_w * (score / 100), Inches(0.05))
        _rect(slide, bar_start, card_y + Inches(0.3), fill_w, Inches(0.2), acc)


def _slide_content_themes(prs, analytics, slide_num):
    slide = _add_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    themes = analytics.content_themes
    companies = list(themes.keys())
    _header(slide, prs, "Content Theme Intelligence",
            "Video topic distribution — what each brand produces and where gaps exist", slide_num)
    all_theme_keys = list(set(k for t in themes.values() for k in t.keys()))
    theme_colors = {th: COLORS_RGB[i % len(COLORS_RGB)] for i, th in enumerate(all_theme_keys)}
    col_w = (W - Inches(0.6)) / max(len(companies), 1)
    for i, c in enumerate(companies[:5]):
        cx = Inches(0.3) + i * col_w
        t_data = themes.get(c, {})
        total = sum(t_data.values()) or 1
        acc = COLORS_RGB[i % len(COLORS_RGB)]
        _rect(slide, cx, Inches(1.1), col_w - Inches(0.1), Inches(0.35), BG_CARD2)
        _rect(slide, cx, Inches(1.1), col_w - Inches(0.1), Inches(0.04), acc)
        _text(slide, c[:14], cx + Inches(0.08), Inches(1.14),
              col_w - Inches(0.2), Inches(0.28), size=9, bold=True, color=acc)
        sorted_themes = sorted(t_data.items(), key=lambda x: x[1], reverse=True)
        for j, (theme, count) in enumerate(sorted_themes[:7]):
            pct = count / total
            row_y = Inches(1.55) + j * Inches(0.74)
            _text(slide, theme[:14], cx + Inches(0.05), row_y,
                  col_w - Inches(0.1), Inches(0.25), size=8, color=TEXT_MUTED)
            bw = col_w - Inches(0.25)
            _rect(slide, cx + Inches(0.05), row_y + Inches(0.27), bw, Inches(0.18),
                  RGBColor(0x1E, 0x29, 0x3B))
            fill = max(bw * pct, Inches(0.05))
            _rect(slide, cx + Inches(0.05), row_y + Inches(0.27), fill, Inches(0.18),
                  theme_colors.get(theme, ACCENT_PURPLE))
            _text(slide, f"{int(pct*100)}%",
                  cx + Inches(0.05) + fill + Inches(0.04),
                  row_y + Inches(0.24), Inches(0.45), Inches(0.22),
                  size=7.5, bold=True, color=TEXT_PRIMARY)


def _slide_heatmap(prs, analytics, user_company, slide_num):
    slide = _add_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    whitespace = analytics.whitespace
    _header(slide, prs, "White Space Opportunity Heatmap",
            "Strategic content gap analysis -- color intensity = opportunity magnitude", slide_num)
    _rect(slide, Inches(0.3), Inches(1.12), W - Inches(0.6), Inches(0.32), BG_CARD2)
    for i, (label, color) in enumerate([("Low (1-3)", ACCENT_RED),
                                         ("Medium (4-6)", ACCENT_AMBER),
                                         ("High (7-10)", ACCENT_GREEN)]):
        lx = Inches(1.5) + i * Inches(3.5)
        _rect(slide, lx, Inches(1.18), Inches(0.18), Inches(0.18), color)
        _text(slide, label, lx + Inches(0.22), Inches(1.17), Inches(1.4), Inches(0.22),
              size=8, color=TEXT_MUTED)
    headers = ["CONTENT THEME", "COVERAGE", "OPP SCORE", "PRIORITY", "YOU?", "STRATEGIC ACTION"]
    col_widths = [Inches(1.9), Inches(1.2), Inches(0.9), Inches(0.95), Inches(0.7), Inches(6.9)]
    col_x = [Inches(0.3)]
    for cw in col_widths[:-1]:
        col_x.append(col_x[-1] + cw)
    header_y = Inches(1.52)
    _rect(slide, Inches(0.3), header_y, W - Inches(0.6), Inches(0.3), RGBColor(0x1E, 0x29, 0x3B))
    for j, (hdr, cx) in enumerate(zip(headers, col_x)):
        _text(slide, hdr, cx + Inches(0.06), header_y + Inches(0.05),
              col_widths[j], Inches(0.22), size=7, bold=True, color=ACCENT_VIOLET)
    for i, row in enumerate(whitespace[:10]):
        row_y = Inches(1.85) + i * Inches(0.51)
        score = row["opportunity_score"]
        cell_color = _heatmap_color(score)
        row_bg = BG_CARD2 if i % 2 == 0 else BG_DARK
        _rect(slide, Inches(0.3), row_y, W - Inches(0.6), Inches(0.48), row_bg)
        priority = "URGENT" if score >= 8 else "HIGH" if score >= 6 else "MEDIUM" if score >= 4 else "LOW"
        pri_color = ACCENT_RED if score >= 8 else ACCENT_AMBER if score >= 6 else ACCENT_GREEN if score >= 4 else TEXT_DIM
        row_vals = [
            (row["theme"],                  TEXT_PRIMARY, False),
            (row["competitor_coverage"],     TEXT_MUTED,   False),
            (f"{score}/10",                  WHITE,        True),
            (priority,                       pri_color,    True),
            ("Yes" if row["you_have_it"] else "No",
             ACCENT_GREEN if row["you_have_it"] else ACCENT_RED, True),
            (row["recommendation"][:72],     TEXT_MUTED,   False),
        ]
        for j, (val, color, bold) in enumerate(row_vals):
            if j == 2:
                _rect(slide, col_x[j] + Inches(0.04), row_y + Inches(0.08),
                      col_widths[j] - Inches(0.08), Inches(0.32), cell_color)
                _text(slide, val, col_x[j] + Inches(0.04), row_y + Inches(0.1),
                      col_widths[j], Inches(0.3), size=9, bold=True,
                      color=WHITE, align=PP_ALIGN.CENTER)
            else:
                _text(slide, val, col_x[j] + Inches(0.06), row_y + Inches(0.1),
                      col_widths[j] - Inches(0.08), Inches(0.32),
                      size=8.5, bold=bold, color=color)


def _slide_threat_matrix(prs, analytics, slide_num):
    slide = _add_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    threat = analytics.threat_scores
    dims = ["Engagement Power", "Audience Loyalty", "Posting Consistency",
            "SEO Optimization", "Content Diversity", "Growth Potential"]
    companies_sorted = sorted(threat.keys(), key=lambda c: threat[c]["total"], reverse=True)
    _header(slide, prs, "Competitive Threat Matrix",
            "Multi-dimensional brand power scoring -- each dimension scored 0-10", slide_num)
    _section_label(slide, "Final Threat Score", Inches(0.3), Inches(1.1))
    for rank, c in enumerate(companies_sorted):
        total = threat[c]["total"]
        acc = COLORS_RGB[rank % len(COLORS_RGB)]
        cy = Inches(1.38) + rank * Inches(0.88)
        _rect(slide, Inches(0.3), cy, Inches(5.5), Inches(0.78), BG_CARD2)
        _rect(slide, Inches(0.3), cy, Inches(0.05), Inches(0.78), acc)
        medal = [" 1 ", " 2 ", " 3 ", " 4 ", " 5 "][rank] if rank < 5 else f"{rank+1}"
        _text(slide, f"{medal}  {c[:16]}", Inches(0.45), cy + Inches(0.06),
              Inches(2.8), Inches(0.3), size=11, bold=True, color=TEXT_PRIMARY)
        _text(slide, f"{total}/10", Inches(4.5), cy + Inches(0.06),
              Inches(0.9), Inches(0.4), size=20, bold=True, color=acc)
        bw = Inches(3.5) * (total / 10)
        _rect(slide, Inches(0.45), cy + Inches(0.5), Inches(3.5), Inches(0.12),
              RGBColor(0x1E, 0x29, 0x3B))
        _rect(slide, Inches(0.45), cy + Inches(0.5), bw, Inches(0.12), acc)
    _section_label(slide, "Dimension Breakdown", Inches(6.2), Inches(1.1))
    grid_x = Inches(6.2)
    col_w = (W - grid_x - Inches(0.3)) / max(len(companies_sorted), 1)
    for i, c in enumerate(companies_sorted):
        acc = COLORS_RGB[i % len(COLORS_RGB)]
        _text(slide, c[:9], grid_x + i * col_w + Inches(0.05), Inches(1.38),
              col_w, Inches(0.3), size=8, bold=True, color=acc, align=PP_ALIGN.CENTER)
    for r, dim in enumerate(dims):
        row_y = Inches(1.72) + r * Inches(0.88)
        _rect(slide, grid_x - Inches(0.05), row_y, W - grid_x + Inches(0.05), Inches(0.82),
              BG_CARD2 if r % 2 == 0 else BG_DARK)
        _text(slide, dim, grid_x, row_y + Inches(0.06),
              Inches(1.8), Inches(0.28), size=8, color=TEXT_MUTED)
        for i, c in enumerate(companies_sorted):
            score = threat[c].get(dim, 0)
            cell_color = _heatmap_color(score)
            cx = grid_x + i * col_w + col_w * 0.1
            cw_cell = col_w * 0.8
            _rect(slide, cx, row_y + Inches(0.38), cw_cell, Inches(0.28), cell_color)
            _text(slide, str(score), cx, row_y + Inches(0.38),
                  cw_cell, Inches(0.28), size=8.5, bold=True,
                  color=WHITE, align=PP_ALIGN.CENTER)


def _slide_recommendations(prs, analytics, user_company, slide_num):
    slide = _add_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    recs = analytics.recommendations
    _header(slide, prs, "AI Strategic Recommendations",
            f"Intelligence-driven strategy for {user_company}", slide_num)
    items = [
        ("1", "Content Strategy",        "content_strategy",  ACCENT_PURPLE),
        ("2", "Upload Cadence",           "upload_cadence",    ACCENT_CYAN),
        ("3", "Shorts Strategy",          "shorts_strategy",   ACCENT_PINK),
        ("4", "Engagement Acceleration",  "engagement_tips",   ACCENT_GREEN),
        ("5", "Differentiation Play",     "differentiation",   ACCENT_AMBER),
    ]
    for i, (num, title, key, acc) in enumerate(items):
        text = recs.get(key, "")[:175]
        ry = Inches(1.12) + i * Inches(1.2)
        _rect(slide, Inches(0.3), ry, W - Inches(0.6), Inches(1.1), BG_CARD2)
        _rect(slide, Inches(0.3), ry, Inches(0.06), Inches(1.1), acc)
        _rect(slide, Inches(0.5), ry + Inches(0.28), Inches(0.38), Inches(0.38), acc)
        _text(slide, num, Inches(0.5), ry + Inches(0.25), Inches(0.4), Inches(0.42),
              size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(slide, title, Inches(1.05), ry + Inches(0.06), Inches(3), Inches(0.3),
              size=10, bold=True, color=acc)
        _text(slide, text, Inches(1.05), ry + Inches(0.38), W - Inches(1.45), Inches(0.65),
              size=9, color=TEXT_PRIMARY)


def _slide_roadmap(prs, analytics, user_company, slide_num):
    slide = _add_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    roadmap = analytics.roadmap
    _header(slide, prs, "30-Day Content Execution Roadmap",
            f"Week-by-week action plan for {user_company}", slide_num)
    week_colors = [ACCENT_PURPLE, ACCENT_PINK, ACCENT_CYAN, ACCENT_GREEN]
    col_w = (W - Inches(0.7)) / 4
    for i, (week, data) in enumerate(roadmap.items()):
        acc = week_colors[i]
        cx = Inches(0.3) + i * (col_w + Inches(0.03))
        _rect(slide, cx, Inches(1.1), col_w, H - Inches(1.45), BG_CARD2)
        _rect(slide, cx, Inches(1.1), col_w, Inches(0.06), acc)
        _rect(slide, cx, Inches(1.16), col_w, Inches(0.38), acc)
        _text(slide, week, cx + Inches(0.1), Inches(1.18), col_w - Inches(0.15), Inches(0.28),
              size=11, bold=True, color=WHITE)
        _text(slide, data["title"], cx + Inches(0.1), Inches(1.6), col_w - Inches(0.15), Inches(0.45),
              size=9, bold=True, color=TEXT_PRIMARY)
        for j, action in enumerate(data["actions"]):
            ay = Inches(2.12) + j * Inches(0.88)
            _rect(slide, cx + Inches(0.1), ay, Inches(0.04), Inches(0.65), acc)
            _text(slide, action[:65], cx + Inches(0.2), ay + Inches(0.04),
                  col_w - Inches(0.3), Inches(0.75), size=8, color=TEXT_PRIMARY)
        _rect(slide, cx, H - Inches(0.68), col_w, Inches(0.38), RGBColor(0x0A, 0x10, 0x1F))
        _text(slide, f"Goal: {data['goal'][:55]}", cx + Inches(0.1), H - Inches(0.66),
              col_w - Inches(0.15), Inches(0.32), size=7.5, italic=True, color=acc)


def _slide_next_steps(prs, analytics, user_company, slide_num):
    slide = _add_slide(prs)
    W, H = prs.slide_width, prs.slide_height
    _header(slide, prs, "Strategic Next Steps",
            "Priority actions to execute within the next 30 days", slide_num)
    top_opp = analytics.whitespace[0]["theme"] if analytics.whitespace else "Tutorial"
    user_freq = analytics.metrics.get(user_company, {}).get("upload_freq", 4)
    steps = [
        (ACCENT_PURPLE, "01", "Claim the Content White Space",
         f"Launch a {top_opp} content series this week -- lowest competitor density, highest organic discovery potential."),
        (ACCENT_PINK,   "02", "Raise Upload Velocity",
         f"Increase to {int(user_freq)+2} videos/month. Build an editorial calendar with fixed weekly publish slots."),
        (ACCENT_CYAN,   "03", "Activate Shorts Engine",
         "Publish 3 Shorts this week. Repurpose top-performing long-form into 60-second hooks with strong opening frames."),
        (ACCENT_GREEN,  "04", "Competitive Monitoring Loop",
         "Audit top 5 videos of your strongest rival weekly. Track title patterns, thumbnail styles, and comment sentiment."),
        (ACCENT_AMBER,  "05", "Engagement Acceleration",
         "Pin a community question within 1 hour of every upload. Reply to all comments in the first 4 hours post-publish."),
    ]
    for i, (acc, num, title, body) in enumerate(steps):
        sy = Inches(1.15) + i * Inches(1.12)
        _rect(slide, Inches(0.3), sy, W - Inches(0.6), Inches(1.02), BG_CARD2)
        _rect(slide, Inches(0.3), sy, Inches(0.06), Inches(1.02), acc)
        _rect(slide, Inches(0.5), sy + Inches(0.25), Inches(0.5), Inches(0.5), acc)
        _text(slide, num, Inches(0.5), sy + Inches(0.23), Inches(0.5), Inches(0.5),
              size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text(slide, title, Inches(1.15), sy + Inches(0.06), Inches(3.5), Inches(0.3),
              size=11, bold=True, color=acc)
        _text(slide, body, Inches(1.15), sy + Inches(0.4), W - Inches(1.6), Inches(0.55),
              size=9.5, color=TEXT_PRIMARY)
    _rect(slide, 0, H - Inches(0.4), W, Inches(0.4), BG_CARD)
    _text(slide, f"AI Video Marketing War Room  .  {user_company}  .  Generated {datetime.now().strftime('%B %Y')}",
          Inches(0.4), H - Inches(0.35), W - Inches(0.8), Inches(0.28),
          size=7.5, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ─── Main Entry Point ─────────────────────────────────────────────────────────

def generate_ppt_report(analytics, channel_data: dict, user_company: str) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    companies = list(analytics.metrics.keys())
    _slide_title(prs, user_company, companies)
    _slide_exec_summary(prs, analytics, user_company, 2)
    _slide_channel_intel(prs, analytics, channel_data, 3)
    _slide_engagement(prs, analytics, 4)
    _slide_aei(prs, analytics, 5)
    _slide_content_themes(prs, analytics, 6)
    _slide_heatmap(prs, analytics, user_company, 7)
    _slide_threat_matrix(prs, analytics, 8)
    _slide_recommendations(prs, analytics, user_company, 9)
    _slide_roadmap(prs, analytics, user_company, 10)
    _slide_next_steps(prs, analytics, user_company, 11)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
